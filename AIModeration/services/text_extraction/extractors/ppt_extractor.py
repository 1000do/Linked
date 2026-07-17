import logging
import io
import time
from typing import Tuple, List, Dict, Any, AsyncIterator
from pptx import Presentation

from core.exceptions import FileProcessingException
from .base_extractor import ITextExtractor

logger = logging.getLogger(__name__)

class PowerPointTextExtractor(ITextExtractor):
    """Extract text from PowerPoint presentation."""

    async def extract_legacy(self, content: bytes) -> Tuple[List[str], float, Dict[str, Any]]:
        if not content:
            raise FileProcessingException("powerpoint", "Empty PPTX provided")
        
        try:
            start_time = time.time()
            prs = Presentation(io.BytesIO(content))
            
            all_text = []
            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    all_text.append(" ".join(slide_texts))
            
            confidence = 1.0 if all_text else 0.0
            processing_time = time.time() - start_time
            
            log_entry = {
                "source": "powerpoint_native",
                "success": True,
                "slides": len(prs.slides),
                "text_length": sum(len(t) for t in all_text),
                "confidence": confidence,
                "processing_time": processing_time
            }
            
            return all_text, confidence, log_entry
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            raise FileProcessingException("powerpoint", f"Extraction failed: {e}")

    async def extract_stream(self, content: bytes) -> AsyncIterator[Tuple[str, float, Dict[str, Any]]]:
        if not content:
            raise FileProcessingException("powerpoint", "Empty PPTX provided")
            
        try:
            prs = Presentation(io.BytesIO(content))
            total_slides = len(prs.slides)
            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                text = " ".join(slide_texts).strip()
                if text:
                    logger.info(f"Extracted native text from PowerPoint slide {slide_idx + 1}: {text[:100]}...")
                    yield text, 1.0, {
                        "source": "powerpoint_native_slide",
                        "success": True,
                        "slide_index": slide_idx + 1,
                        "total_slides": total_slides,
                        "text_length": len(text),
                    }
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            raise FileProcessingException("powerpoint", f"Extraction failed: {e}")
