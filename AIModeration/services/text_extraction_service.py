"""Service for extracting text from various media types."""

import logging
import io
import tempfile
import os
import math
import fitz
from typing import Tuple, AsyncIterator, Dict, Any, Optional, List
from PIL import Image
import cv2
import easyocr
import torch
from difflib import SequenceMatcher
from faster_whisper import WhisperModel
from pdf2image import convert_from_bytes
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import numpy as np
import warnings

# Suppress annoying PyTorch DataLoader warnings on CPU
warnings.filterwarnings("ignore", category=UserWarning, module="torch.utils.data.dataloader")

from core.exceptions import FileProcessingException
from services.base_service import BaseService
from config.settings import get_settings, Settings

logger = logging.getLogger(__name__)


class TextExtractionService(BaseService):
    """Extract text from media files using OCR and Whisper."""
    
    # Class-level model cache (singleton pattern)
    _whisper_model = None
    _ocr_reader = None
    _models_loaded = False
    
    def __init__(self, settings: Optional[Settings] = None):
        """Initialize text extraction service."""
        super().__init__("TextExtractionService")
        self.settings = settings or get_settings()
        
        if not TextExtractionService._models_loaded:
            self._load_models(self.settings)
            
    @classmethod
    def _load_models(cls, settings: Settings):
        """Load models once (class-level cache)."""
        if cls._models_loaded:
            return
            
        logger.info("Loading text extraction models...")
        
        try:
            whisper_model_name = settings.WHISPER_MODEL_NAME
            
            device = settings.DEVICE
            compute_type = "float16" if torch.cuda.is_available() else "int8"
            
            logger.info(f"Loading Whisper model {whisper_model_name} on {device}...")
            cls._whisper_model = WhisperModel(whisper_model_name, device=device, compute_type=compute_type)
            
            logger.info("Loading EasyOCR model...")
            cls._ocr_reader = easyocr.Reader(['en','vi'], gpu=torch.cuda.is_available(), verbose=False)
            
            cls._models_loaded = True
            logger.info("✓ All text extraction models loaded successfully")
            
        except Exception as e:
            logger.error(f"Failed to load text extraction models: {e}")

    def get_file_type_for_text_extraction(self, file_extension: str) -> str:
        """
        Map file extension to appropriate generic text extraction file type.
        """
        if not file_extension:
            return "text"
        ext = file_extension.lower().replace(".", "").strip()
        if ext in ["txt", "text","csv"]:
            return "text"
        elif ext in ["jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp","image"]:
            return "image"
        elif ext in ["mp4", "avi", "mov", "mkv", "webm", "flv", "wmv", "3gp","video"]:
            return "video"
        elif ext in ["pdf"]:
            return "pdf"
        elif ext in ["docx", "doc", "odt", "rtf","word"]:
            return "word"
        elif ext in ["pptx", "ppt","powerpoint"]:
            return "powerpoint"
        elif ext in ["xlsx", "xls","excel"]:
            return "excel"
        return "text"
    
    async def extract_from_image(self, image_bytes: bytes) -> Tuple[List[str], float, dict]:
        """
        Extract text from image using OCR.
        """
        if not image_bytes:
            raise FileProcessingException("image", "Empty image provided")
        
        try:
            with self.time_operation("extract_from_image", size=len(image_bytes)):
                # Load image
                image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                
                # Run OCR
                # EasyOCR expects a numpy array, byte content, or filepath
                # Since image is loaded as PIL, we convert to numpy array
                image_np = np.array(image)
                results = self._ocr_reader.readtext(image_np, detail=1, paragraph=False)
                
                texts = []
                confidences = []
                for res in results:
                    if len(res) == 3:
                        _, t, prob = res
                        if t.strip():
                            texts.append(t)
                            confidences.append(float(prob))
                            
                text = texts
                confidence = sum(confidences) / len(confidences) if confidences else 0.0
                
                log_entry = {
                    "source": "image_ocr",
                    "success": True,
                    "text_length": sum(len(t) for t in text),
                    "confidence": confidence,
                }
                
                logger.info(f"Extracted OCR text from image: {len(text)} segments")
                return text, confidence, log_entry
        
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            raise FileProcessingException("image", f"OCR failed: {e}")
    
    async def extract_from_pdf_legacy(self, pdf_bytes: bytes) -> Tuple[List[str], float, dict]:
        """
        Extract text from PDF using native extraction with OCR fallback.
        """
        if not pdf_bytes:
            raise FileProcessingException("pdf", "Empty PDF provided")
        
        try:
            with self.time_operation("extract_from_pdf", size=len(pdf_bytes)):
                # Try native extraction first
                doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                total_pages = len(doc)
                all_text = []
                has_native_text = False
                
                for idx in range(total_pages):
                    page = doc.load_page(idx)
                    text = page.get_text().strip()
                    if text:
                        has_native_text = True
                        all_text.append(f"Page {idx+1}: {text}")
                
                doc.close()

                if has_native_text:
                    confidence = 1.0
                    log_entry = {
                        "source": "pdf_native",
                        "success": True,
                        "pages_processed": total_pages,
                        "text_length": sum(len(t) for t in all_text),
                        "confidence": confidence,
                    }
                    return all_text, confidence, log_entry

                logger.info("No native text found in PDF, falling back to OCR")
                # Fallback to OCR
                images = convert_from_bytes(pdf_bytes, first_page=1, last_page=5)  # First 5 pages
                
                all_text = []
                all_confidences = []
                
                for idx, image in enumerate(images):
                    image_np = np.array(image)
                    results = self._ocr_reader.readtext(image_np, detail=1, paragraph=False)
                    page_texts = []
                    for res in results:
                        if len(res) == 3:
                            _, t, prob = res
                            if t.strip():
                                page_texts.append(t)
                                all_confidences.append(float(prob))
                    if page_texts:
                        all_text.append(f"Page {idx+1}: " + " ".join(page_texts))
                
                confidence = sum(all_confidences) / len(all_confidences) if all_confidences else 0.0
                
                log_entry = {
                    "source": "pdf_ocr",
                    "success": True,
                    "pages_processed": len(images),
                    "text_length": sum(len(t) for t in all_text),
                    "confidence": confidence,
                }
                
                return all_text, confidence, log_entry
        
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise FileProcessingException("pdf", f"Extraction failed: {e}")
    
    async def extract_from_video_legacy(
        self,
        video_bytes: bytes,
        sample_interval_seconds: float = 3.0,
        max_frames: int = 30,
    ) -> Tuple[List[str], float, dict]:
        """
        Extract text from video frames and audio using OCR and Whisper.
        """
        if not video_bytes:
            raise FileProcessingException("video", "Empty video provided")
        
        try:
            with self.time_operation("extract_from_video", size=len(video_bytes)):
                # Write to temp file
                with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as tmp:
                    tmp.write(video_bytes)
                    tmp_path = tmp.name
                
                all_text = []
                all_confidences = []
                
                try:
                    # Detect scenes using PySceneDetect
                    import scenedetect
                    from scenedetect import ContentDetector
                    scene_list = scenedetect.detect(tmp_path, ContentDetector(threshold=27.0), show_progress=True)
                    logger.info(f"Scenes detected: {len(scene_list)}")
                    cap = cv2.VideoCapture(tmp_path)
                    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
                    frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                    
                    target_frames = []
                    for scene in scene_list:
                        start_frame = scene[0].frame_num
                        end_frame = scene[1].frame_num
                        middle_frame = start_frame + (end_frame - start_frame) // 2
                        target_frames.append(middle_frame)
                        
                    if len(target_frames) > max_frames:
                        # Pick the longest scenes
                        scenes_with_duration = [(scene[1].frame_num - scene[0].frame_num, idx) for idx, scene in enumerate(scene_list)]
                        scenes_with_duration.sort(reverse=True, key=lambda x: x[0])
                        top_indices = [x[1] for x in scenes_with_duration[:max_frames]]
                        top_indices.sort() # Keep chronological order
                        target_frames = [target_frames[i] for i in top_indices]
                    elif len(target_frames) == 0:
                        # Fallback if no scenes detected
                        total_duration = frame_count / fps if fps > 0 else 0
                        interval = max(sample_interval_seconds, total_duration / max_frames if max_frames > 0 else sample_interval_seconds)
                        sample_every_n_frames = max(1, int(fps * interval))
                        target_frames = list(range(0, frame_count, sample_every_n_frames))[:max_frames]
                    
                    frames_processed = 0
                    last_text = ""
                    
                    for frame_idx in target_frames:
                        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
                        ret, frame = cap.read()
                        if not ret:
                            continue
                            
                        rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                        
                        # Downscale for faster OCR
                        height, width = rgb_frame.shape[:2]
                        max_width = 800
                        if width > max_width:
                            scale = max_width / width
                            new_width = max_width
                            new_height = int(height * scale)
                            rgb_frame = cv2.resize(rgb_frame, (new_width, new_height))
                        
                        results = self._ocr_reader.readtext(rgb_frame, detail=1, paragraph=False)
                        frame_texts = []
                        frame_confs = []
                        for res in results:
                            if len(res) == 3:
                                _, t, prob = res
                                if t.strip():
                                    frame_texts.append(t)
                                    frame_confs.append(float(prob))
                                    
                        text = " ".join(frame_texts).strip()
                        frame_conf = sum(frame_confs) / len(frame_confs) if frame_confs else 0.0
                        
                        if text:
                            # Deduplication logic
                            similarity = SequenceMatcher(None, last_text, text).ratio()
                            if similarity < 0.9: # Only add if it's less than 80% similar to previous frame
                                # all_text.append(f"Video Frame OCR {frame_idx}: {text}")
                                all_text.append(text)
                                all_confidences.append(frame_conf)
                                last_text = text
                        
                        frames_processed += 1
                    
                    cap.release()
                except Exception as e:
                    logger.error(f"OCR text extraction failed: {e}")
                
                # Extract audio using Whisper (faster-whisper for efficiency)
                try:
                    # Save audio track
                    audio_path = tmp_path.replace(".mp4", ".wav")
                    # Extract as 16kHz mono WAV for optimal Whisper quality
                    os.system(f"ffmpeg -i {tmp_path} -vn -acodec pcm_s16le -ar 16000 -ac 1 -y {audio_path} 2>/dev/null")
                    
                    if os.path.exists(audio_path) and self._whisper_model:
                        segments, _ = self._whisper_model.transcribe(
                            audio_path,
                            beam_size=5,
                            condition_on_previous_text=False
                        )
                        
                        logger.info(f'Audio text extracted using whisper')
                        for seg_idx, segment in enumerate(segments):
                            seg_text = segment.text.strip()
                            if seg_text:
                                # all_text.append(f"Video Segment Audio {seg_idx + 1}: {seg_text}")
                                all_text.append(seg_text)
                                all_confidences.append(math.exp(segment.avg_logprob))
                        
                        os.unlink(audio_path)
                
                except Exception as e:
                    logger.warning(f"Audio extraction failed (continuing with OCR): {e}")
                    
                finally:
                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)
                
                confidence = sum(all_confidences) / len(all_confidences) if all_confidences else 0.0
                
                log_entry = {
                    "source": "video_ocr_whisper",
                    "success": True,
                    "total_frames": frame_count,
                    "frames_sampled": frames_processed,
                    "text_length": sum(len(t) for t in all_text),
                    "confidence": confidence,
                }
                
                return all_text, confidence, log_entry
        
        except Exception as e:
            logger.error(f"Video extraction failed: {e}")
            raise FileProcessingException("video", f"Extraction failed: {e}")
    
    async def extract_from_word_legacy(self, docx_bytes: bytes) -> Tuple[List[str], float, dict]:
        """
        Extract text from Word document.
        """
        if not docx_bytes:
            raise FileProcessingException("word", "Empty DOCX provided")
        
        try:
            with self.time_operation("extract_from_word", size=len(docx_bytes)):
                # Load from bytes
                doc = Document(io.BytesIO(docx_bytes))
                
                # Extract all text
                all_text = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        all_text.append(para.text)
                
                confidence = 1.0 if all_text else 0.0  # High confidence for text extraction
                
                log_entry = {
                    "source": "word_native",
                    "success": True,
                    "paragraphs": len(doc.paragraphs),
                    "text_length": sum(len(t) for t in all_text),
                    "confidence": confidence,
                }
                
                return all_text, confidence, log_entry
        
        except Exception as e:
            logger.error(f"Word extraction failed: {e}")
            raise FileProcessingException("word", f"Extraction failed: {e}")
    
    async def extract_from_powerpoint_legacy(self, pptx_bytes: bytes) -> Tuple[List[str], float, dict]:
        """
        Extract text from PowerPoint presentation.
        """
        if not pptx_bytes:
            raise FileProcessingException("powerpoint", "Empty PPTX provided")
        
        try:
            with self.time_operation("extract_from_powerpoint", size=len(pptx_bytes)):
                # Load from bytes
                prs = Presentation(io.BytesIO(pptx_bytes))
                
                # Extract text from slides
                all_text = []
                for slide_idx, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    if slide_texts:
                        all_text.append(f"Slide {slide_idx+1}: " + " ".join(slide_texts))
                
                confidence = 1.0 if all_text else 0.0
                
                log_entry = {
                    "source": "powerpoint_native",
                    "success": True,
                    "slides": len(prs.slides),
                    "text_length": sum(len(t) for t in all_text),
                    "confidence": confidence,
                }
                
                return all_text, confidence, log_entry
        
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            raise FileProcessingException("powerpoint", f"Extraction failed: {e}")
    
    async def extract_from_excel_legacy(self, xlsx_bytes: bytes) -> Tuple[List[str], float, dict]:
        """
        Extract text from Excel spreadsheet.
        """
        if not xlsx_bytes:
            raise FileProcessingException("excel", "Empty XLSX provided")
        
        try:
            with self.time_operation("extract_from_excel", size=len(xlsx_bytes)):
                # Load from bytes
                wb = load_workbook(io.BytesIO(xlsx_bytes))
                
                # Extract all cell values
                all_text = []
                for sheet in wb.sheetnames[:5]:  # First 5 sheets
                    ws = wb[sheet]
                    sheet_texts = []
                    for row in ws.iter_rows(values_only=True):
                        for cell in row:
                            if cell and str(cell).strip():
                                sheet_texts.append(str(cell).strip())
                    if sheet_texts:
                        all_text.append(f"Sheet {sheet}: " + " ".join(sheet_texts))
                
                confidence = 1.0 if all_text else 0.0
                
                log_entry = {
                    "source": "excel_native",
                    "success": True,
                    "sheets_processed": min(5, len(wb.sheetnames)),
                    "text_length": sum(len(t) for t in all_text),
                    "confidence": confidence,
                }
                
                return all_text, confidence, log_entry
        
        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            raise FileProcessingException("excel", f"Extraction failed: {e}")
    
    async def extract_generic_legacy(
        self,
        content: bytes,
        material_type: str,
    ) -> Tuple[List[str], float, dict]:
        """
        Extract text from generic material type.
        """
        material_type = material_type.lower().strip()
        
        if material_type in ["text"]:
            try:
                text = content.decode("utf-8")
                return [text], 1.0, {"source": "text_native", "success": True}
            except UnicodeDecodeError:
                raise FileProcessingException("text", "Invalid UTF-8 encoding")
        
        elif material_type in ["image"]:
            return await self.extract_from_image(content)
        
        elif material_type in ["pdf"]:
            return await self.extract_from_pdf_legacy(content)
        
        elif material_type in ["video"]:
            return await self.extract_from_video_legacy(content)
        
        elif material_type in ["word"]:
            return await self.extract_from_word_legacy(content)
        
        elif material_type in ["powerpoint"]:
            return await self.extract_from_powerpoint_legacy(content)
        
        elif material_type in ["excel"]:
            return await self.extract_from_excel_legacy(content)
        
        else:
            raise FileProcessingException(material_type, f"Unsupported file type: {material_type}")

    async def extract_from_video_stream(
        self,
        video_bytes: bytes,
        sample_interval_seconds: float = 3.0,
        max_frames: int = 30,
    ) -> AsyncIterator[Tuple[str, float, Dict[str, Any]]]:
        """
        Extract text from video frames (OCR) and audio (Whisper) as a stream of segments.
        Yields Tuple[text, confidence, details_dict].
        """
        if not video_bytes:
            raise FileProcessingException("video", "Empty video provided")
        
        tmp_path = None
        audio_path = None
        
        try:
            with self.time_operation("extract_from_video_stream", size=len(video_bytes)):
                # Write to temp file
                with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as tmp:
                    tmp.write(video_bytes)
                    tmp_path = tmp.name
                
                # 1. OCR text extraction
                try:
                    import scenedetect
                    from scenedetect import ContentDetector
                    scene_list = scenedetect.detect(tmp_path, ContentDetector(threshold=27.0))
                    
                    cap = cv2.VideoCapture(tmp_path)
                    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
                    frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                    
                    target_frames = []
                    for scene in scene_list:
                        start_frame = scene[0].frame_num
                        end_frame = scene[1].frame_num
                        middle_frame = start_frame + (end_frame - start_frame) // 2
                        target_frames.append(middle_frame)
                        
                    if len(target_frames) > max_frames:
                        # Pick the longest scenes
                        scenes_with_duration = [(scene[1].frame_num - scene[0].frame_num, idx) for idx, scene in enumerate(scene_list)]
                        scenes_with_duration.sort(reverse=True, key=lambda x: x[0])
                        top_indices = [x[1] for x in scenes_with_duration[:max_frames]]
                        top_indices.sort() # Keep chronological order
                        target_frames = [target_frames[i] for i in top_indices]
                    elif len(target_frames) == 0:
                        total_duration = frame_count / fps if fps > 0 else 0
                        interval = max(sample_interval_seconds, total_duration / max_frames if max_frames > 0 else sample_interval_seconds)
                        sample_every_n_frames = max(1, int(fps * interval))
                        target_frames = list(range(0, frame_count, sample_every_n_frames))[:max_frames]
                    
                    frames_processed = 0
                    last_text = ""
                    last_frame = 0
                    
                    for frame_idx in target_frames:
                        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
                        ret, frame = cap.read()
                        if not ret:
                            continue
                        
                        rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                        
                        # Downscale for faster OCR
                        height, width = rgb_frame.shape[:2]
                        max_width = 800
                        if width > max_width:
                            scale = max_width / width
                            new_width = max_width
                            new_height = int(height * scale)
                            rgb_frame = cv2.resize(rgb_frame, (new_width, new_height))
                        
                        results = self._ocr_reader.readtext(rgb_frame, detail=1, paragraph=False)
                        
                        frame_texts = []
                        frame_confs = []
                        for res in results:
                            if len(res) == 3:
                                _, t, prob = res
                                if t.strip():
                                    frame_texts.append(t)
                                    frame_confs.append(float(prob))
                                    
                        text = " | ".join(frame_texts).strip()
                        frame_conf = sum(frame_confs) / len(frame_confs) if frame_confs else 0.0
                        
                        if text:
                            # Deduplication logic
                            similarity = SequenceMatcher(None, last_text, text).ratio()
                            if similarity < 0.8:
                                logger.info(f"OCR similarity score of frame {frame_idx} vs frame {last_frame}: {similarity}")
                                last_text = text
                                last_frame = frame_idx
                                logger.info(f'OCR on video frame {frame_idx}: {text}')
                                yield text, frame_conf, {
                                    "source": "video_ocr_frame",
                                    "success": True,
                                    "frame_index": frame_idx,
                                    "total_frames": frame_count,
                                    "text_length": len(text),
                                }
                                
                        frames_processed += 1
                        
                    cap.release()
                except Exception as e:
                    logger.error(f"OCR text extraction failed: {e}")
                
                # 2. Whisper audio text extraction
                try:
                    # Save audio track
                    audio_path = tmp_path.replace(".mp4", ".wav")
                    os.system(f"ffmpeg -i {tmp_path} -vn -acodec pcm_s16le -ar 16000 -ac 1 -y {audio_path} 2>/dev/null")
                    
                    if os.path.exists(audio_path) and self._whisper_model:
                        segments, _ = self._whisper_model.transcribe(
                            audio_path,
                            beam_size=5,
                            condition_on_previous_text=False
                        )
                        
                        for idx,segment in enumerate(segments):
                            audio_text = segment.text
                            duration = segment.end - segment.start
                            segment_conf = math.exp(segment.avg_logprob)
                            logger.info(f'Audio transcription on Segment {idx+1} (Duration: {duration:.2f}s, [{segment.start:.2f}s - {segment.end:.2f}s]): {audio_text}')
                            if audio_text.strip():
                                yield segment.text, segment_conf, {
                                    "source": "video_whisper_segment",
                                    "success": True,
                                    "segment_start": segment.start,
                                    "segment_end": segment.end,
                                    "text_length": len(segment.text),
                                }
                
                except Exception as e:
                    logger.warning(f"Audio extraction failed (continuing with OCR only): {e}")
        
        except Exception as e:
            logger.error(f"Video extraction failed: {e}")
            raise FileProcessingException("video", f"Extraction failed: {e}")
            
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except Exception as ex:
                    logger.warning(f"Failed to delete temp video file {tmp_path}: {ex}")
            if audio_path and os.path.exists(audio_path):
                try:
                    os.unlink(audio_path)
                except Exception as ex:
                    logger.warning(f"Failed to delete temp audio file {audio_path}: {ex}")

    
    async def extract_from_pdf_stream(self, pdf_bytes: bytes) -> AsyncIterator[Tuple[str, float, Dict[str, Any]]]:
        if not pdf_bytes:
            raise FileProcessingException("pdf", "Empty PDF provided")
        try:
            with self.time_operation("extract_from_pdf_stream", size=len(pdf_bytes)):
                doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                total_pages = len(doc)
                has_native_text = False
                
                for idx in range(total_pages):
                    page = doc.load_page(idx)
                    text = page.get_text().strip()
                    if text:
                        has_native_text = True
                        logger.info(f"Extracted native text from PDF page {idx + 1}: {text[:100]}...")
                        yield text, 1.0, {
                            "source": "pdf_native_page",
                            "success": True,
                            "page": idx + 1,
                            "total_pages": total_pages,
                            "text_length": len(text),
                        }
                
                doc.close()
                
                if not has_native_text:
                    logger.info("No native text found in PDF, falling back to OCR")
                    images = convert_from_bytes(pdf_bytes)
                    total_pages = len(images)
                    for idx, image in enumerate(images):
                        image_np = np.array(image)
                        results = self._ocr_reader.readtext(image_np, detail=1, paragraph=False)
                        page_texts = []
                        page_confs = []
                        for res in results:
                            if len(res) == 3:
                                _, t, prob = res
                                if t.strip():
                                    page_texts.append(t)
                                    page_confs.append(float(prob))
                        text = " | ".join(page_texts).strip()
                        confidence = sum(page_confs) / len(page_confs) if page_confs else 0.0
                        if text:
                            logger.info(f"Extracted OCR text from PDF page {idx + 1}: {text[:100]}...")
                            yield text, confidence, {
                                "source": "pdf_ocr_page",
                                "success": True,
                                "page": idx + 1,
                                "total_pages": total_pages,
                                "text_length": len(text),
                            }
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise FileProcessingException("pdf", f"Extraction failed: {e}")

    async def extract_from_word_stream(self, docx_bytes: bytes) -> AsyncIterator[Tuple[str, float, Dict[str, Any]]]:
        if not docx_bytes:
            raise FileProcessingException("word", "Empty DOCX provided")
        try:
            with self.time_operation("extract_from_word_stream", size=len(docx_bytes)):
                doc = Document(io.BytesIO(docx_bytes))
                total_paras = len(doc.paragraphs)
                for idx, para in enumerate(doc.paragraphs):
                    text = para.text.strip()
                    if text:
                        logger.info(f"Extracted native text from Word paragraph {idx + 1}: {text[:100]}...")
                        yield text, 1.0, {
                            "source": "word_native_paragraph",
                            "success": True,
                            "paragraph_index": idx + 1,
                            "total_paragraphs": total_paras,
                            "text_length": len(text),
                        }
        except Exception as e:
            logger.error(f"Word extraction failed: {e}")
            raise FileProcessingException("word", f"Extraction failed: {e}")

    async def extract_from_powerpoint_stream(self, pptx_bytes: bytes) -> AsyncIterator[Tuple[str, float, Dict[str, Any]]]:
        if not pptx_bytes:
            raise FileProcessingException("powerpoint", "Empty PPTX provided")
        try:
            with self.time_operation("extract_from_powerpoint_stream", size=len(pptx_bytes)):
                prs = Presentation(io.BytesIO(pptx_bytes))
                total_slides = len(prs.slides)
                for slide_idx, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    text = " | ".join(slide_texts).strip()
                    if text:
                        logger.info(f"Extracted native text from PowerPoint slide {slide_idx + 1}: {text[:100]}...")
                        yield text, 1.0, {
                            "source": "powerpoint_native_slide",
                            "success": True,
                            "slide_index": slide_idx + 1,
                            "total_slides": total_slides,
                            "text_length": len(text),
                        }
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            raise FileProcessingException("powerpoint", f"Extraction failed: {e}")

    async def extract_from_excel_stream(self, xlsx_bytes: bytes) -> AsyncIterator[Tuple[str, float, Dict[str, Any]]]:
        if not xlsx_bytes:
            raise FileProcessingException("excel", "Empty XLSX provided")
        try:
            with self.time_operation("extract_from_excel_stream", size=len(xlsx_bytes)):
                wb = load_workbook(io.BytesIO(xlsx_bytes))
                total_sheets = len(wb.sheetnames)
                for sheet_idx, sheet in enumerate(wb.sheetnames):
                    ws = wb[sheet]
                    sheet_texts = []
                    for row in ws.iter_rows(values_only=True):
                        for cell in row:
                            if cell and str(cell).strip():
                                sheet_texts.append(str(cell).strip())
                    text = " ".join(sheet_texts).strip()
                    if text:
                        logger.info(f"Extracted native text from Excel sheet '{sheet}': {text[:100]}...")
                        yield text, 1.0, {
                            "source": "excel_native_sheet",
                            "success": True,
                            "sheet_name": sheet,
                            "sheet_index": sheet_idx + 1,
                            "total_sheets": total_sheets,
                            "text_length": len(text),
                        }
        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            raise FileProcessingException("excel", f"Extraction failed: {e}")

    async def extract_generic_stream(
        self,
        content: bytes,
        material_type: str,
    ) -> AsyncIterator[Tuple[str, float, Dict[str, Any]]]:
        """
        Extract text from generic material type as a stream of segments.
        Yields Tuple[text, confidence, details_dict].
        """
        material_type = material_type.lower().strip()
        
        if material_type in ["text"]:
            try:
                text = content.decode("utf-8")
                logger.info(f"Extracted text from generic text: {text[:100]}...")
                yield text, 1.0, {"source": "text_native", "success": True}
            except UnicodeDecodeError:
                raise FileProcessingException("text", "Invalid UTF-8 encoding")
        
        elif material_type in ["image"]:
            text, conf, log = await self.extract_from_image(content)
            yield text, conf, log
            
        elif material_type in ["pdf"]:
            async for text, conf, log in self.extract_from_pdf_stream(content):
                yield text, conf, log
            
        elif material_type in ["video"]:
            async for text, conf, log in self.extract_from_video_stream(content, sample_interval_seconds=3.0):
                yield text, conf, log
                
        elif material_type in ["word"]:
            async for text, conf, log in self.extract_from_word_stream(content):
                yield text, conf, log
            
        elif material_type in ["powerpoint"]:
            async for text, conf, log in self.extract_from_powerpoint_stream(content):
                yield text, conf, log
            
        elif material_type in ["excel"]:
            async for text, conf, log in self.extract_from_excel_stream(content):
                yield text, conf, log
            
        else:
            raise FileProcessingException(material_type, f"Unsupported file type: {material_type}")
